# services/document_service.py
"""
Document Processing Service

This module provides a comprehensive document processing service that can extract text content
from various document formats and prepare them for text analysis, search, and AI processing.

WHAT THIS FILE DOES:
====================
1. **Multi-format Document Support**: Extracts text from PDF, DOCX, Excel, CSV, TXT, Markdown, 
   HTML, EPUB, RTF, ODT, and PowerPoint files.

2. **Intelligent Text Extraction**: Uses multiple extraction methods with fallback strategies
   to ensure reliable text extraction even from complex or corrupted documents.

3. **Smart Content Chunking**: Breaks down large documents into semantically meaningful chunks
   that preserve context while being optimal for search and AI processing.

4. **Rich Metadata Extraction**: Extracts and generates comprehensive metadata including content
   analysis, performance metrics, and document statistics.

5. **Performance Monitoring**: Tracks processing times, extraction rates, and resource usage
   for optimization and monitoring purposes.

MAIN WORKFLOW:
==============
Input Document → Text Extraction → Content Preprocessing → Semantic Chunking → Metadata Generation → Output

KEY FEATURES:
=============
- Supports 11+ document formats
- Fallback extraction strategies for reliability  
- Semantic content analysis and classification
- Performance optimization for large files
- Comprehensive logging and error handling
- Memory-efficient processing with progress tracking

TYPICAL USE CASES:
==================
- Document indexing for search systems
- Preparing documents for RAG (Retrieval Augmented Generation) systems
- Content analysis and classification
- Document processing pipelines
- Knowledge base construction

DEPENDENCIES:
=============
- PyMuPDF/fitz: PDF processing
- pdfplumber: Alternative PDF extraction
- python-docx: Word document processing
- openpyxl: Excel file processing
- BeautifulSoup: HTML parsing
- ebooklib: EPUB processing
- LangChain: Document chunking and management
"""

import fitz
import pdfplumber
import os
import hashlib
import time
from datetime import datetime
from typing import List, Dict, Any, Optional
import re
import numpy as np
from collections import Counter
import logging

import docx
import openpyxl
import csv
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
from rtfparse.parser import Rtf_Parser
from striprtf.striprtf import rtf_to_text
from odf import text, teletype
from odf.opendocument import load
from pptx import Presentation
from services.ocr_processor import OCRProcessor
from PIL import Image

from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

class DocumentService:
    """
    A comprehensive service for processing various document formats and extracting text content.
    
    This service handles the complete pipeline from document ingestion to chunked output,
    including text extraction, preprocessing, semantic chunking, and metadata generation.
    
    Attributes:
        chunk_size (int): Maximum size of text chunks in characters
        chunk_overlap (int): Overlap between adjacent chunks for context preservation
        text_splitter (RecursiveCharacterTextSplitter): LangChain text splitter instance
    
    Supported Formats:
        - PDF (.pdf) - Using PyMuPDF and pdfplumber
        - Word Documents (.docx) - Using python-docx
        - Excel Files (.xlsx, .xlsm) - Using openpyxl
        - CSV Files (.csv) - Using built-in csv module
        - Text Files (.txt) - Plain text reading
        - Markdown (.md) - Plain text reading with markdown preservation
        - HTML (.html) - Using BeautifulSoup for clean text extraction
        - EPUB (.epub) - Using ebooklib for ebook processing
        - RTF (.rtf) - Using rtfparse for rich text format
        - OpenDocument Text (.odt) - Using odfpy
        - PowerPoint (.pptx, .ppt) - Using python-pptx
    """
    
    def __init__(self, chunk_size: int = 300, chunk_overlap: int = 50):
        """
        Initialize the DocumentService with specified chunking parameters.
        
        Args:
            chunk_size (int): Maximum number of characters per chunk. Default is 300.
            chunk_overlap (int): Number of characters to overlap between chunks for 
                                context preservation. Default is 50.
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            # Hierarchical separators: prefer paragraph breaks, then sentences, then clauses
            separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""]
        )

        self.ocr_processor = OCRProcessor()
    
    def _extract_from_image(self, filepath: str) -> Dict[str, Any]:
        """Extract text from image files using OCR."""
        try:
            if not self.ocr_processor.is_available():
                return {
                    'success': False,
                    'error': 'OCR processor not available. Please install easyocr: pip install easyocr'
                }
            
            ocr_result = self.ocr_processor.process_image_file(filepath)
            
            if not ocr_result['success']:
                return {
                    'success': False,
                    'error': f"OCR failed: {ocr_result.get('error', 'Unknown error')}"
                }
            
            text_content = ocr_result['text']
            if text_content and text_content.strip():
                metadata = {
                    'page': 1,
                    'source': filepath,
                    'extraction_method': 'easyocr',
                    'content_type': 'image_ocr',
                    'ocr_confidence': ocr_result['confidence'],
                    'word_count': ocr_result.get('word_count', 0),
                    'image_size': ocr_result.get('image_size', 'unknown'),
                    'processed_at': ocr_result.get('processed_at')
                }
                return {
                    'success': True,
                    'documents': [Document(page_content=text_content, metadata=metadata)],
                    'method': 'easyocr'
                }
            
            return {
                'success': False,
                'error': 'No text found in image'
            }
            
        except Exception as e:
            logger.error(f"Image OCR extraction failed: {str(e)}")
            return {
                'success': False,
                'error': f"Failed to extract text from image: {str(e)}"
            }

    def process_document(self, filepath: str, file_id: str, original_filename: str) -> Dict[str, Any]:
        """
        Main entry point for document processing. Orchestrates the complete processing pipeline.
        
        This method coordinates text extraction, chunking, and metadata generation while
        providing comprehensive timing and performance metrics.
        
        Args:
            filepath (str): Full path to the document file to be processed
            file_id (str): Unique identifier for the document (used in chunk IDs)
            original_filename (str): Original name of the file (may differ from current filename)
        
        Returns:
            Dict[str, Any]: Processing result containing:
                - success (bool): Whether processing completed successfully
                - chunks (List[Document]): List of processed document chunks (if successful)
                - metadata (Dict): Document metadata and statistics (if successful)
                - processing_stats (Dict): Performance metrics and timing data (if successful)
                - error (str): Error message (if unsuccessful)
        
        Processing Pipeline:
            1. File validation and initial metadata collection
            2. Text extraction using format-specific methods
            3. Content preprocessing and cleaning
            4. Semantic chunking with overlap management
            5. Metadata extraction and enrichment
            6. Performance metrics calculation
        """
        try:
            # Validate file existence and readability
            if not os.path.exists(filepath):
                return {
                    'success': False,
                    'error': f"File does not exist: {filepath}"
                }

            if not os.access(filepath, os.R_OK):
                return {
                    'success': False,
                    'error': f"File is not readable: {filepath}"
                }
            start_time = time.time()
            file_size = os.path.getsize(filepath) if os.path.exists(filepath) else 0
            file_ext = os.path.splitext(original_filename)[1].lower()
            
            logger.info(f"Starting document processing: {original_filename} "
                       f"({file_ext}, {file_size:,} bytes)")
            
            # Track text extraction time
            extraction_start = time.time()
            logger.info(f"Starting text extraction for {original_filename}")
            
            # Extract text content using format-specific methods
            text_extraction_result = self._extract_text_content(filepath)
            
            # Ensure we have a valid result dictionary
            if text_extraction_result is None:
                return {
                    'success': False,
                    'error': "Text extraction returned None - internal error"
                }
            
            extraction_time = time.time() - extraction_start
            logger.info(f"Text extraction completed for {original_filename}: "
                       f"{extraction_time:.2f}s")
            
            # Check if extraction was successful
            if not text_extraction_result['success']:
                return {
                    'success': False,
                    'error': text_extraction_result['error']
                }
            
            documents = text_extraction_result['documents']
            total_chars = sum(len(doc.page_content) for doc in documents)
            
            logger.info(f"Extracted {len(documents)} pages, {total_chars:,} characters "
                       f"from {original_filename}")
            
            # Track chunking time
            chunking_start = time.time()
            logger.info(f"Starting chunking for {original_filename}")
            
            # Create semantic chunks for optimal retrieval
            chunks = self._create_semantic_chunks(documents, file_id, original_filename)
            
            chunking_time = time.time() - chunking_start
            logger.info(f"Chunking completed for {original_filename}: "
                       f"{chunking_time:.2f}s, {len(chunks)} chunks created")
            
            # Track metadata extraction time
            metadata_start = time.time()
            
            # Extract comprehensive document metadata
            metadata = self._extract_document_metadata(filepath, original_filename, documents)
            
            metadata_time = time.time() - metadata_start
            processing_time = time.time() - start_time
            
            logger.info(f"Document processing completed for {original_filename}: "
                       f"Total: {processing_time:.2f}s, "
                       f"Extraction: {extraction_time:.2f}s, "
                       f"Chunking: {chunking_time:.2f}s, "
                       f"Metadata: {metadata_time:.2f}s")
            
            # Generate comprehensive processing statistics
            processing_stats = {
                'total_pages': len(documents),
                'total_chunks': len(chunks),
                'total_characters': total_chars,
                'total_words': sum(len(doc.page_content.split()) for doc in documents),
                'processing_method': text_extraction_result['method'],
                'processed_at': datetime.now().isoformat(),
                'file_size': file_size,
                'file_extension': file_ext,
                'timing': {
                    'total_processing_time': processing_time,
                    'extraction_time': extraction_time,
                    'chunking_time': chunking_time,
                    'metadata_time': metadata_time,
                    'extraction_rate_chars_per_sec': total_chars / extraction_time if extraction_time > 0 else 0,
                    'processing_rate_mb_per_sec': (file_size / 1024 / 1024) / processing_time if processing_time > 0 else 0,
                    'chunking_rate_chars_per_sec': total_chars / chunking_time if chunking_time > 0 else 0
                }
            }
            
            # Log performance warnings for slow processing
            if file_ext in ['.xlsx', '.xlsm']:
                logger.warning(f"Excel file processing time for {original_filename}: "
                              f"{processing_time:.2f}s ({processing_stats['timing']['processing_rate_mb_per_sec']:.2f} MB/s)")
            
            logger.info(f"Document processed successfully: {len(chunks)} chunks created")
            
            return {
                'success': True,
                'chunks': chunks,
                'metadata': metadata,
                'processing_stats': processing_stats
            }
            
        except Exception as e:
            processing_time = time.time() - start_time if 'start_time' in locals() else 0
            logger.error(f"Document processing error for {original_filename} "
                        f"after {processing_time:.2f}s: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }
    
    def _extract_text_content(self, filepath: str) -> Dict[str, Any]:
        """
        Route document to appropriate extraction method based on file extension.
        
        This method acts as a dispatcher, selecting the optimal extraction method
        for each supported file format. Each format has specialized handling to
        maximize extraction quality and handle format-specific quirks.
        
        Args:
            filepath (str): Path to the document file
        
        Returns:
            Dict[str, Any]: Extraction result containing:
                - success (bool): Whether extraction succeeded
                - documents (List[Document]): Extracted document objects (if successful)
                - method (str): Extraction method used (if successful)
                - error (str): Error description (if unsuccessful)
        
        Supported Formats and Methods:
            - .pdf: PyMuPDF (primary) → pdfplumber (fallback)
            - .docx: python-docx library
            - .xlsx/.xlsm: openpyxl with progress tracking
            - .csv: Built-in csv module
            - .txt: Direct file reading
            - .md: Markdown-aware text reading
            - .html: BeautifulSoup with tag filtering
            - .epub: ebooklib with chapter extraction
            - .rtf: rtfparse for rich text
            - .odt: odfpy for OpenDocument
            - .pptx/.ppt: python-pptx for presentations
        """
        file_ext = os.path.splitext(filepath)[1].lower()
        
        # Format dispatcher with comprehensive format support
        if file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
            return self._extract_from_image(filepath)
        elif file_ext == '.pdf':
            return self._extract_from_pdf(filepath)
        # ... rest of your existing format handlers
        elif file_ext == '.docx':
            return self._extract_from_docx(filepath)
        elif file_ext == '.xlsx' or file_ext == '.xlsm':
            return self._extract_from_excel(filepath)
        elif file_ext == '.csv':
            return self._extract_from_csv(filepath)
        elif file_ext == '.txt':
            return self._extract_from_txt(filepath)
        elif file_ext == '.md':
            return self._extract_from_md(filepath)
        elif file_ext == '.html':
            return self._extract_from_html(filepath)
        elif file_ext == '.epub':
            return self._extract_from_epub(filepath)
        elif file_ext == '.rtf':
            return self._extract_from_rtf(filepath)
        elif file_ext == '.odt':
            return self._extract_from_odt(filepath)
        elif file_ext == '.pptx' or file_ext == '.ppt':
            return self._extract_from_pptx(filepath)
        else:
            return {
                'success': False,
                'error': f"Unsupported file type: {file_ext}"
            }

    # ==========================================
    # FORMAT-SPECIFIC EXTRACTION METHODS
    # ==========================================
    
    def _extract_from_pdf(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from PDF files using multiple extraction strategies.
        
        Uses a two-tier approach:
        1. PyMuPDF (fitz) - Fast and handles most PDFs well
        2. pdfplumber - Fallback for complex layouts and tables
        
        Args:
            filepath (str): Path to PDF file
        
        Returns:
            Dict[str, Any]: Extraction result with documents and method used
        """
        documents = []
        
        # Primary extraction method: PyMuPDF
        try:
            pdf_doc = fitz.open(filepath)
            
            try:
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc.load_page(page_num)
                    text = page.get_text()
                    
                    if text and text.strip():
                        # Digital text found
                        metadata = {
                            'page': page_num + 1,
                            'source': filepath,
                            'extraction_method': 'pymupdf',
                            'content_type': 'digital_text'
                        }
                        documents.append(Document(page_content=text, metadata=metadata))
                    else:
                        # No digital text, try OCR
                        if self.ocr_processor.is_available():
                            logger.info(f"No digital text on page {page_num + 1}, attempting OCR...")
                            ocr_result = self.ocr_processor.perform_ocr_on_pdf_page(page)
                            
                            if ocr_result['success'] and ocr_result['text'].strip():
                                metadata = {
                                    'page': page_num + 1,
                                    'source': filepath,
                                    'extraction_method': 'pymupdf_ocr_fallback',
                                    'content_type': 'scanned_text',
                                    'ocr_confidence': ocr_result['confidence']
                                }
                                documents.append(Document(page_content=ocr_result['text'], metadata=metadata))
                            else:
                                logger.warning(f"OCR failed for page {page_num + 1}: {ocr_result.get('error', 'No text found')}")
                        else:
                            logger.warning(f"No text found on page {page_num + 1} and OCR not available")
            finally:
                pdf_doc.close()
                
            if documents:
                return {'success': True, 'documents': documents, 'method': 'pymupdf_with_ocr'}
        except Exception as e:
            logger.warning(f"PyMuPDF extraction failed: {str(e)}")

        # Fallback extraction method: pdfplumber (without OCR)
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        metadata = {
                            'page': page_num + 1,
                            'source': filepath,
                            'extraction_method': 'pdfplumber',
                            'content_type': 'digital_text'
                        }
                        documents.append(Document(page_content=text, metadata=metadata))
                        
                if documents:
                    return {'success': True, 'documents': documents, 'method': 'pdfplumber'}
        except Exception as e:
            logger.error(f"pdfplumber extraction failed: {str(e)}")
        
        # If we reach here, all extraction methods failed
        return {
            'success': False,
            'error': "Failed to extract text from PDF using all available methods"
        }
    
    def get_ocr_stats(self) -> Dict[str, Any]:
        """Get OCR processor statistics."""
        return self.ocr_processor.get_ocr_stats()
        
    def _extract_from_docx(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from Microsoft Word documents (.docx).
        
        Processes all paragraphs in the document and combines them into
        a single text block with proper paragraph separation.
        
        Args:
            filepath (str): Path to DOCX file
        
        Returns:
            Dict[str, Any]: Extraction result with document content
        """
        try:
            doc = docx.Document(filepath)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            text_content = "\n\n".join(full_text)
            
            if text_content and text_content.strip():
                metadata = {'page': 1, 'source': filepath, 'extraction_method': 'python-docx'}
                return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'python-docx'}
            
            return {'success': False, 'error': "No text found in DOCX file"}
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from DOCX: {str(e)}"}
            
    def _extract_from_pptx(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from PowerPoint presentation files (.pptx, .ppt).
        
        Processes each slide separately and extracts text from all text frames
        and shapes within the slide.
        
        Args:
            filepath (str): Path to PowerPoint file
        
        Returns:
            Dict[str, Any]: Extraction result with documents (one per slide)
        """
        try:
            prs = Presentation(filepath)
            documents = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            slide_text.append("".join([run.text for run in paragraph.runs]))
                
                text_content = "\n".join(slide_text)
                if text_content.strip():
                    metadata = {'page': slide_num + 1, 'source': filepath, 'extraction_method': 'python-pptx'}
                    documents.append(Document(page_content=text_content, metadata=metadata))
            
            if documents:
                return {'success': True, 'documents': documents, 'method': 'python-pptx'}
            
            return {'success': False, 'error': "No text found in PPTX file"}
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from PPTX: {str(e)}"}
    
    # ==========================================
    # TEXT PROCESSING AND CHUNKING METHODS
    # ==========================================
    
    def _create_semantic_chunks(self, documents: List[Document], file_id: str, filename: str) -> List[Document]:
        """
        Create semantically meaningful chunks from extracted documents.
        
        This method implements advanced chunking strategies that preserve context
        and meaning while creating optimal sizes for retrieval and processing.
        
        Key Features:
        - Reduced chunk size (800 chars) for better retrieval precision
        - Increased overlap (400 chars) for better context preservation  
        - Semantic scoring to assess chunk quality
        - Content type detection for categorization
        - Key term extraction for enhanced searchability
        - Comprehensive metadata enrichment
        
        Args:
            documents (List[Document]): List of extracted document pages/sections
            file_id (str): Unique identifier for the source document
            filename (str): Original filename for reference
        
        Returns:
            List[Document]: List of processed chunks with enhanced metadata
        
        Chunk Metadata Includes:
            - file_id: Source document identifier
            - filename: Original document name
            - chunk_id: Unique chunk identifier
            - chunk_index: Position within document
            - semantic_score: Quality/richness score (0.0-1.0)
            - content_density: Words per character ratio
            - content_type: Detected content category
            - key_terms: Extracted important terms
            - chunk_created_at: Processing timestamp
        """
        chunks = []
        
        for doc in documents:
            # Clean and preprocess content for better chunking
            processed_content = self._preprocess_content(doc.page_content)
            
            # Enhanced text splitter with optimized parameters
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""]
            )
            
            try:
                base_chunks = text_splitter.split_text(processed_content)
            except Exception as e:
                logger.warning(f"Text splitting failed, falling back to simple chunking: {str(e)}")
                # Fallback to simple chunking
                chunk_size = self.chunk_size
                base_chunks = [processed_content[i:i+chunk_size] 
                            for i in range(0, len(processed_content), chunk_size)]
            
            for i, chunk_content in enumerate(base_chunks):
                # Skip very small chunks (likely noise)
                if len(chunk_content.strip()) < 50:
                    continue
                
                # Create enhanced metadata for each chunk
                enhanced_metadata = doc.metadata.copy()
                enhanced_metadata.update({
                    'file_id': file_id,
                    'filename': filename,
                    'chunk_id': f"{file_id}_{doc.metadata.get('page', 0)}_{i}",
                    'chunk_index': i,
                    'semantic_score': self._calculate_semantic_score(chunk_content),
                    'content_density': len(chunk_content.split()) / len(chunk_content),
                    'content_type': self._detect_content_type(chunk_content),
                    'key_terms': ', '.join(self._extract_key_terms(chunk_content)[:5]),
                    'chunk_created_at': datetime.now().isoformat(),
                })
                
                chunks.append(Document(
                    page_content=chunk_content,
                    metadata=enhanced_metadata
                ))
        
        return chunks
    
    def _preprocess_content(self, content: str) -> str:
        """
        Preprocess text content to improve chunking quality.
        
        Applies various text normalization and cleanup operations to handle
        common issues from document extraction, particularly PDF artifacts.
        
        Operations:
        - Whitespace normalization
        - Missing space insertion (common in PDF extraction)
        - Sentence boundary correction
        
        Args:
            content (str): Raw extracted text content
        
        Returns:
            str: Cleaned and normalized text content
        """
        # Normalize excessive whitespace
        content = re.sub(r'\s+', ' ', content)
        
        # Fix common PDF extraction issues
        content = re.sub(r'([a-z])([A-Z])', r'\1 \2', content)  # Fix missing spaces between words
        content = re.sub(r'([.!?])([A-Z])', r'\1 \2', content)  # Fix sentence boundaries
        
        return content.strip()
    
    def _calculate_semantic_score(self, content: str) -> float:
        """
        Calculate a semantic richness score for content quality assessment.
        
        This score helps identify high-quality, information-rich chunks that
        are likely to be useful for retrieval and analysis.
        
        Scoring Factors:
        - Word diversity (unique words / total words) - 40% weight
        - Presence of verbs (indicates actionable content) - 30% weight  
        - Presence of entities (proper nouns, names) - 20% weight
        - Presence of numbers (data, facts, figures) - 10% weight
        
        Args:
            content (str): Text content to score
        
        Returns:
            float: Semantic richness score between 0.0 and 1.0
        """
        words = content.split()
        if not words:
            return 0.0
        
        # Calculate word diversity ratio
        unique_words = len(set(words))
        word_diversity = unique_words / len(words)
        
        # Check for content richness indicators
        has_verbs = bool(re.search(r'\b(is|are|was|were|have|has|do|does|will|would|can|could)\b', content.lower()))
        has_entities = bool(re.search(r'\b[A-Z][a-z]+\b', content))  # Proper nouns
        has_numbers = bool(re.search(r'\d+', content))  # Numerical data
        
        # Weighted scoring formula
        score = (word_diversity * 0.4 +
                (1 if has_verbs else 0) * 0.3 +
                (1 if has_entities else 0) * 0.2 +
                (1 if has_numbers else 0) * 0.1)
        
        return min(1.0, score)  # Cap at 1.0
    
    def _detect_content_type(self, text: str) -> str:
        """
        Detect the type/category of content based on keyword patterns.
        
        This classification helps with content organization and can be used
        for specialized processing or retrieval strategies.
        
        Content Types Detected:
        - introduction: Introductory or background content
        - methodology: Methods, procedures, approaches
        - results: Findings, data, analysis results
        - conclusion: Summaries, conclusions, discussions
        - abstract: Document abstracts or summaries
        - references: Citations, bibliographies
        - technical: Technical content with algorithms, formulas
        - content: General content (default category)
        
        Args:
            text (str): Text content to classify
        
        Returns:
            str: Detected content type category
        """
        text_lower = text.lower()
        
        # Define keyword patterns for each content type
        type_patterns = {
            'introduction': ['introduction', 'background', 'overview'],
            'methodology': ['method', 'methodology', 'approach', 'procedure'],
            'results': ['result', 'finding', 'outcome', 'data', 'analysis'],
            'conclusion': ['conclusion', 'discussion', 'summary'],
            'abstract': ['abstract', 'summary'],
            'references': ['reference', 'bibliography', 'citation'],
            'technical': ['algorithm', 'formula', 'equation', 'implementation'],
        }
        
        # Check for pattern matches
        for content_type, keywords in type_patterns.items():
            if any(keyword in text_lower for keyword in keywords):
                return content_type
        
        return 'content'  # Default category
    
    def _extract_key_terms(self, text: str) -> List[str]:
        """
        Extract key terms and important words from text content.
        
        Identifies the most significant terms in the content by frequency
        analysis while filtering out common stop words and short words.
        
        Process:
        1. Extract words (3-15 characters, alphabetic only)
        2. Convert to lowercase for consistency
        3. Filter out common stop words
        4. Count word frequency
        5. Return top 10 most frequent terms
        
        Args:
            text (str): Text content to analyze
        
        Returns:
            List[str]: List of key terms ordered by frequency
        """
        # Extract meaningful words (3-15 chars, alphabetic)
        words = re.findall(r'\b[a-zA-Z]{3,15}\b', text.lower())
        
        # Common stop words to filter out
        stop_words = {
            'the', 'and', 'for', 'with', 'this', 'that', 'from', 'are', 'was', 'were',
            'which', 'have', 'been', 'will', 'can', 'may', 'also', 'such', 'these'
        }
        
        # Filter words and count frequency
        filtered_words = [word for word in words if word not in stop_words and len(word) > 3]
        word_counts = Counter(filtered_words)
        
        return [word for word, _ in word_counts.most_common(10)]
    
    def _extract_document_metadata(self, filepath: str, filename: str, documents: List[Document]) -> Dict[str, Any]:
        """
        Extract comprehensive metadata from processed documents.
        
        Generates detailed statistics and metadata about the document content,
        including size metrics, content analysis, and processing information.
        
        Metadata Categories:
        - Basic Info: filename, size, page count
        - Content Stats: word/character counts, averages
        - Content Analysis: content types, key terms
        - Processing Info: timestamps, language detection
        
        Args:
            filepath (str): Path to source document file
            filename (str): Original document filename  
            documents (List[Document]): Processed document pages/sections
        
        Returns:
            Dict[str, Any]: Comprehensive document metadata
        """
        # Calculate content statistics
        total_words = sum(len(doc.page_content.split()) for doc in documents)
        total_chars = sum(len(doc.page_content) for doc in documents)
        
        # Analyze content composition
        content_types = Counter(self._detect_content_type(doc.page_content) for doc in documents)
        
        # Extract key terms from sample content (first 5 pages)
        sample_text = " ".join([doc.page_content for doc in documents[:5]])
        key_terms = self._extract_key_terms(sample_text)
        
        return {
            # Basic document information
            'filename': filename,
            'file_size': os.path.getsize(filepath) if os.path.exists(filepath) else 0,
            'total_pages': len(documents),
            
            # Content statistics  
            'total_words': total_words,
            'total_characters': total_chars,
            'average_words_per_page': total_words / max(len(documents), 1),
            
            # Content analysis
            'content_types': dict(content_types.most_common()),
            'key_terms': key_terms[:20],  # Top 20 key terms
            
            # Processing metadata
            'processed_at': datetime.now().isoformat(),
            'language': 'english'  # Could be enhanced with automatic language detection
        }
    
    def _process_excel_sheet_efficiently(self, sheet, sheet_name):
        """Process Excel sheet with memory optimization for large files."""
        text_content = ""
        batch_size = 1000  # Process in batches
        current_batch = []
        
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                current_batch.append(row_text)
                
                if len(current_batch) >= batch_size:
                    text_content += "\n".join(current_batch) + "\n"
                    current_batch = []
        
        # Process remaining rows
        if current_batch:
            text_content += "\n".join(current_batch)
        
        return text_content

    def _extract_from_excel(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from Excel files (.xlsx, .xlsm) with performance monitoring.
        
        Processes each worksheet separately and converts cell data to tab-separated text.
        Includes comprehensive progress tracking for large files and performance optimization.
        
        Features:
        - Progress tracking for large worksheets
        - Per-sheet timing and statistics
        - Memory-efficient row-by-row processing
        - Performance warnings for slow processing
        
        Args:
            filepath (str): Path to Excel file
        
        Returns:
            Dict[str, Any]: Extraction result with documents (one per sheet)
        """
        try:
            start_time = time.time()
            file_size = os.path.getsize(filepath)
            logger.info(f"Starting Excel extraction: {os.path.basename(filepath)} ({file_size:,} bytes)")
            
            # Track workbook loading time
            load_start = time.time()
            try:
                # Add read-only mode for better performance and memory usage
                workbook = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
            except Exception as e:
                # Fallback to regular mode if read-only fails
                logger.warning(f"Read-only mode failed, falling back to regular mode: {str(e)}")
                workbook = openpyxl.load_workbook(filepath, data_only=True)
            load_time = time.time() - load_start
            
            sheet_count = len(workbook.sheetnames)
            logger.info(f"Excel workbook loaded in {load_time:.2f}s, {sheet_count} sheets found")
            
            documents = []
            total_cells_processed = 0
            
            # Process each worksheet
            for i, sheet_name in enumerate(workbook.sheetnames):
                sheet_start = time.time()
                logger.info(f"Processing sheet {i+1}/{sheet_count}: {sheet_name}")
                
                sheet = workbook[sheet_name]
                text_content = ""
                cells_in_sheet = 0
                
                # Get sheet dimensions for progress tracking
                max_row = sheet.max_row
                max_col = sheet.max_column
                total_cells_in_sheet = max_row * max_col
                
                logger.info(f"Sheet {sheet_name} dimensions: {max_row} rows x {max_col} columns "
                           f"({total_cells_in_sheet:,} cells)")
                
                # Process rows with progress tracking
                for row_num, row in enumerate(sheet.iter_rows(), 1):
                    row_data = []
                    for cell in row:
                        cells_in_sheet += 1
                        if cell.value is not None:
                            row_data.append(str(cell.value))
                        else:
                            row_data.append("")
                    
                    # Only add non-empty rows
                    if any(cell_data.strip() for cell_data in row_data):
                        text_content += "\t".join(row_data) + "\n"
                    
                    # Progress logging for large sheets
                    if max_row > 1000 and row_num % 1000 == 0:
                        logger.info(f"Processed {row_num}/{max_row} rows in sheet {sheet_name}")
                
                sheet_time = time.time() - sheet_start
                total_cells_processed += cells_in_sheet
                
                logger.info(f"Sheet {sheet_name} processed in {sheet_time:.2f}s: "
                           f"{cells_in_sheet:,} cells, {len(text_content):,} characters")
                
                # Create document for each sheet with detailed metadata
                if text_content.strip():
                    metadata = {
                        'page': sheet_name, 
                        'source': filepath, 
                        'extraction_method': 'openpyxl',
                        'sheet_index': i,
                        'rows': max_row,
                        'columns': max_col,
                        'cells_processed': cells_in_sheet,
                        'processing_time': sheet_time
                    }
                    documents.append(Document(page_content=text_content, metadata=metadata))
            
            extraction_time = time.time() - start_time
            
            logger.info(f"Excel extraction completed in {extraction_time:.2f}s: "
                       f"{len(documents)} sheets, {total_cells_processed:,} total cells, "
                       f"{sum(len(doc.page_content) for doc in documents):,} characters")
            
            # Performance metrics calculation
            cells_per_sec = total_cells_processed / extraction_time if extraction_time > 0 else 0
            mb_per_sec = (file_size / 1024 / 1024) / extraction_time if extraction_time > 0 else 0
            
            logger.info(f"Excel processing rate: {cells_per_sec:.0f} cells/sec, {mb_per_sec:.2f} MB/sec")
            
            # Performance warning for slow processing
            if extraction_time > 10:
                logger.warning(f"Slow Excel processing detected: {extraction_time:.2f}s for {file_size:,} bytes "
                              f"({sheet_count} sheets, {total_cells_processed:,} cells)")
            
            if documents:
                return {'success': True, 'documents': documents, 'method': 'openpyxl'}
            
            return {'success': False, 'error': "No text found in Excel file"}
            
        except Exception as e:
            extraction_time = time.time() - start_time if 'start_time' in locals() else 0
            logger.error(f"Excel extraction failed after {extraction_time:.2f}s: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from Excel: {str(e)}"}

    def _extract_from_csv(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from CSV files.
        
        Reads CSV data and converts to comma-separated text format,
        preserving the tabular structure in text form.
        
        Args:
            filepath (str): Path to CSV file
        
        Returns:
            Dict[str, Any]: Extraction result with CSV content
        """
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                text_content = "\n".join([",".join(row) for row in reader])
            
            if text_content and text_content.strip():
                metadata = {'page': 1, 'source': filepath, 'extraction_method': 'csv'}
                return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'csv'}
            
            return {'success': False, 'error': "No text found in CSV file"}
        except Exception as e:
            logger.error(f"CSV extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from CSV: {str(e)}"}
            
    def _extract_from_txt(self, filepath: str) -> Dict[str, Any]:
        """Extract text from plain text files with encoding fallback."""
        encodings_to_try = ['utf-8', 'latin1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings_to_try:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    text_content = f.read()
                
                if text_content and text_content.strip():
                    metadata = {
                        'page': 1, 
                        'source': filepath, 
                        'extraction_method': 'txt',
                        'encoding_used': encoding
                    }
                    return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'txt'}
                
                return {'success': False, 'error': "No text found in TXT file"}
                
            except UnicodeDecodeError:
                logger.warning(f"Failed to decode with {encoding}, trying next encoding")
                continue
            except Exception as e:
                logger.error(f"TXT extraction failed with {encoding}: {str(e)}")
                continue
        
        return {'success': False, 'error': "Failed to decode text file with any supported encoding"}

    def _extract_from_md(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from Markdown files.
        
        Reads Markdown content as plain text, preserving formatting markers.
        Future enhancement could include Markdown parsing for structure.
        
        Args:
            filepath (str): Path to Markdown file
        
        Returns:
            Dict[str, Any]: Extraction result with Markdown content
        """
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            # Note: Currently processed as plain text
            # Future enhancement: Parse markdown for semantic structure
            
            if text_content and text_content.strip():
                metadata = {'page': 1, 'source': filepath, 'extraction_method': 'md'}
                return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'md'}
            
            return {'success': False, 'error': "No text found in Markdown file"}
        except Exception as e:
            logger.error(f"Markdown extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from Markdown: {str(e)}"}

    def _extract_from_html(self, filepath: str) -> Dict[str, Any]:
        """
        Extract clean text from HTML files.
        
        Uses BeautifulSoup to parse HTML and extract clean text content,
        removing scripts, styles, and HTML tags.
        
        Args:
            filepath (str): Path to HTML file
        
        Returns:
            Dict[str, Any]: Extraction result with clean text content
        """
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
            
            # Remove script and style elements for clean text extraction
            for script in soup(["script", "style"]):
                script.extract()
            
            text_content = soup.get_text()
            
            if text_content and text_content.strip():
                metadata = {'page': 1, 'source': filepath, 'extraction_method': 'beautifulsoup'}
                return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'beautifulsoup'}
            
            return {'success': False, 'error': "No text found in HTML file"}
        except Exception as e:
            logger.error(f"HTML extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from HTML: {str(e)}"}
            
    def _extract_from_epub(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from EPUB ebook files.
        
        Processes each chapter/document in the EPUB as a separate document,
        using BeautifulSoup to clean HTML content from chapters.
        
        Args:
            filepath (str): Path to EPUB file
        
        Returns:
            Dict[str, Any]: Extraction result with documents (one per chapter)
        """
        try:
            book = epub.read_epub(filepath)
            documents = []
            
            # Process each document item in the EPUB
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_content = soup.get_text()
                    if text_content.strip():
                        metadata = {'page': item.get_name(), 'source': filepath, 'extraction_method': 'ebooklib'}
                        documents.append(Document(page_content=text_content, metadata=metadata))
            
            if documents:
                return {'success': True, 'documents': documents, 'method': 'ebooklib'}
                
            return {'success': False, 'error': "No text found in EPUB file"}
        except Exception as e:
            logger.error(f"EPUB extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from EPUB: {str(e)}"}
            
    def _extract_from_rtf(self, filepath: str) -> Dict[str, Any]:
        """Extract text from Rich Text Format (RTF) files."""
        try:
            with open(filepath, 'rb') as f:
                rtf_content = f.read()
                text_content = rtf_to_text(rtf_content.decode('utf-8', errors='ignore'))
                    
            if text_content and text_content.strip():
                metadata = {'page': 1, 'source': filepath, 'extraction_method': 'striprtf'}
                return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'striprtf'}
            
            return {'success': False, 'error': "No text found in RTF file"}
        except Exception as e:
            logger.error(f"RTF extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from RTF: {str(e)}"}
            
    def _extract_from_odt(self, filepath: str) -> Dict[str, Any]:
        """
        Extract text from OpenDocument Text (ODT) files.
        
        Uses odfpy library to process ODT files and extract paragraph text.
        
        Args:
            filepath (str): Path to ODT file
        
        Returns:
            Dict[str, Any]: Extraction result with ODT content
        """
        try:
            odt_doc = load(filepath)
            text_content = ""
            # Extract text from all paragraph elements
            for elem in odt_doc.getElementsByType(text.P):
                text_content += teletype.extractText(elem) + "\n"
            if text_content.strip():
                metadata = {'page': 1, 'source': filepath, 'extraction_method': 'odfpy'}
                return {'success': True, 'documents': [Document(page_content=text_content, metadata=metadata)], 'method': 'odfpy'}
            return {'success': False, 'error': "No text found in ODT file"}
        except Exception as e:
            logger.error(f"ODT extraction failed: {str(e)}")
            return {'success': False, 'error': f"Failed to extract text from ODT: {str(e)}"}