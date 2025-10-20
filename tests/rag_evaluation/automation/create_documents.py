"""
Document Generator for RAG Evaluation
Creates documents in various formats (DOCX, PDF, PPTX, CSV) for comprehensive testing
"""

import os
from pathlib import Path
from docx import Document
from docx.shared import Inches
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt
from PIL import Image, ImageDraw, ImageFont
import csv

def create_technical_docx():
    """Create technical documentation in DOCX format"""
    doc = Document()
    
    # Title
    title = doc.add_heading('XploreaseAPI Documentation', 0)
    
    # Overview section
    doc.add_heading('Overview', level=1)
    overview = doc.add_paragraph(
        'The XploreaseAPI is a RESTful web service that provides document processing and '
        'retrieval capabilities for enterprise applications. Version 2.1.0 includes enhanced '
        'security features and improved performance.'
    )
    
    # Installation section
    doc.add_heading('Installation', level=1)
    doc.add_heading('System Requirements', level=2)
    requirements = doc.add_paragraph()
    requirements.add_run('• Python 3.8 or higher\n')
    requirements.add_run('• Memory: Minimum 4GB RAM, Recommended 8GB\n')
    requirements.add_run('• Storage: 2GB free space\n')
    requirements.add_run('• Operating System: Windows 10+, macOS 10.15+, Ubuntu 18.04+')
    
    doc.add_heading('Installation Steps', level=2)
    steps = doc.add_paragraph()
    steps.add_run('1. Download the installation package from https://releases.xplorease.com/v2.1.0\n')
    steps.add_run('2. Extract the archive to your desired directory (default: C:\\Program Files\\Xplorease)\n')
    steps.add_run('3. Run the setup script: python setup.py install\n')
    steps.add_run('4. Configure environment variables:\n')
    steps.add_run('   - XPLOREASE_HOME: Installation directory\n')
    steps.add_run('   - XPLOREASE_CONFIG: Configuration file path\n')
    steps.add_run('5. Start the service: xplorease-server start')
    
    # API Endpoints section
    doc.add_heading('API Endpoints', level=1)
    
    doc.add_heading('Authentication', level=2)
    auth_para = doc.add_paragraph()
    auth_para.add_run('POST /api/v2/auth/login').bold = True
    auth_para.add_run('\n• Description: Authenticate user and obtain access token\n')
    auth_para.add_run('• Parameters: username, password\n')
    auth_para.add_run('• Response: JWT token valid for 24 hours\n')
    auth_para.add_run('• Rate limit: 5 requests per minute')
    
    doc.add_heading('Document Upload', level=2)
    upload_para = doc.add_paragraph()
    upload_para.add_run('POST /api/v2/documents/upload').bold = True
    upload_para.add_run('\n• Description: Upload documents for processing\n')
    upload_para.add_run('• Supported formats: PDF, DOCX, TXT, CSV, JSON\n')
    upload_para.add_run('• Maximum file size: 50MB\n')
    upload_para.add_run('• Concurrent uploads: Maximum 3 files\n')
    upload_para.add_run('• Processing time: 30-120 seconds depending on file size')
    
    doc.add_heading('Document Search', level=2)
    search_para = doc.add_paragraph()
    search_para.add_run('GET /api/v2/documents/search').bold = True
    search_para.add_run('\n• Description: Search through uploaded documents\n')
    search_para.add_run('• Parameters: query (required), limit (optional, default 10), offset (optional, default 0)\n')
    search_para.add_run('• Response time: Typically 200-500ms\n')
    search_para.add_run('• Supports fuzzy matching and semantic search')
    
    # Troubleshooting section
    doc.add_heading('Troubleshooting', level=1)
    doc.add_heading('Common Issues', level=2)
    
    doc.add_heading('Installation Fails', level=3)
    install_issue = doc.add_paragraph()
    install_issue.add_run('Symptom: ').bold = True
    install_issue.add_run('Setup script exits with error code 1\n')
    install_issue.add_run('Cause: ').bold = True
    install_issue.add_run('Insufficient permissions or missing dependencies\n')
    install_issue.add_run('Solution: ').bold = True
    install_issue.add_run('Run as administrator and ensure Python pip is updated')
    
    doc.add_heading('High Memory Usage', level=3)
    memory_issue = doc.add_paragraph()
    memory_issue.add_run('Symptom: ').bold = True
    memory_issue.add_run('Service consumes >6GB RAM\n')
    memory_issue.add_run('Cause: ').bold = True
    memory_issue.add_run('Large document processing or memory leak\n')
    memory_issue.add_run('Solution: ').bold = True
    memory_issue.add_run('Restart service and reduce concurrent upload limit')
    
    doc.add_heading('Slow Query Performance', level=3)
    perf_issue = doc.add_paragraph()
    perf_issue.add_run('Symptom: ').bold = True
    perf_issue.add_run('Search responses take >2 seconds\n')
    perf_issue.add_run('Cause: ').bold = True
    perf_issue.add_run('Database indexing issues or large result sets\n')
    perf_issue.add_run('Solution: ').bold = True
    perf_issue.add_run('Rebuild search indexes using xplorease-admin reindex')
    
    # Performance Metrics
    doc.add_heading('Performance Metrics', level=1)
    perf_para = doc.add_paragraph()
    perf_para.add_run('• Average response time: 150ms for search queries\n')
    perf_para.add_run('• Throughput: 1000 requests per minute\n')
    perf_para.add_run('• Uptime: 99.9% availability target\n')
    perf_para.add_run('• Document processing rate: 50 documents per minute')
    
    # Support
    doc.add_heading('Support', level=1)
    support_para = doc.add_paragraph()
    support_para.add_run('• Documentation: https://docs.xplorease.com\n')
    support_para.add_run('• Support email: support@xplorease.com\n')
    support_para.add_run('• Emergency hotline: +1-800-XPLOREASE\n')
    support_para.add_run('• Business hours: Monday-Friday, 9 AM - 6 PM EST')
    
    return doc

def create_business_pdf():
    """Create business report in PDF format"""
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    
    # Create document
    filename = Path(__file__).parent / "../documents/business/financial_report_q3_2024.pdf"
    doc = SimpleDocTemplate(str(filename), pagesize=letter, topMargin=0.5*inch)
    
    # Get styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Title'],
        fontSize=18,
        spaceAfter=30,
        textColor=colors.darkblue
    )
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading1'],
        fontSize=14,
        spaceAfter=12,
        textColor=colors.darkblue
    )
    
    # Build story
    story = []
    
    # Title
    story.append(Paragraph("Q3 2024 Financial Report", title_style))
    story.append(Paragraph("<b>TechCorp Industries</b>", styles['Normal']))
    story.append(Spacer(1, 20))
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", heading_style))
    story.append(Paragraph(
        "TechCorp Industries delivered exceptional performance in Q3 2024, with revenue growth of 23% "
        "year-over-year and improved operational efficiency across all business units. The company continues "
        "to strengthen its market position in cloud services and artificial intelligence solutions.",
        styles['Normal']
    ))
    story.append(Spacer(1, 12))
    
    # Financial Performance
    story.append(Paragraph("Financial Performance", heading_style))
    
    # Revenue Analysis
    story.append(Paragraph("<b>Revenue Analysis</b>", styles['Heading2']))
    revenue_data = [
        ['Metric', 'Q3 2024', 'Q3 2023', 'Growth'],
        ['Total Revenue', '$485.2M', '$394.8M', '22.9%'],
        ['Quarterly Growth', '$485.2M', '$449.1M (Q2)', '8.1%'],
    ]
    revenue_table = Table(revenue_data, colWidths=[2*inch, 1.5*inch, 1.5*inch, 1*inch])
    revenue_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(revenue_table)
    story.append(Spacer(1, 12))
    
    # Revenue by Segment
    story.append(Paragraph("<b>Revenue by Segment</b>", styles['Heading2']))
    segment_data = [
        ['Segment', 'Revenue', 'Percentage'],
        ['Cloud Services', '$298.4M', '61.5%'],
        ['AI Solutions', '$124.7M', '25.7%'],
        ['Traditional Software', '$62.1M', '12.8%'],
    ]
    segment_table = Table(segment_data, colWidths=[2*inch, 1.5*inch, 1.5*inch])
    segment_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.darkblue),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.lightblue),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(segment_table)
    story.append(Spacer(1, 12))
    
    # Profitability Metrics
    story.append(Paragraph("<b>Profitability Metrics</b>", styles['Heading2']))
    profit_list = [
        "• <b>Gross Profit:</b> $321.8 million (66.3% margin)",
        "• <b>Operating Income:</b> $97.4 million (20.1% margin)",
        "• <b>Net Income:</b> $73.2 million (15.1% margin)",
        "• <b>Earnings Per Share (EPS):</b> $2.18 vs $1.65 in Q3 2023"
    ]
    for item in profit_list:
        story.append(Paragraph(item, styles['Normal']))
    story.append(Spacer(1, 12))
    
    # Customer Acquisition
    story.append(Paragraph("Customer Acquisition", heading_style))
    customer_list = [
        "• <b>New Customers:</b> 2,847 new enterprise clients added",
        "• <b>Customer Retention Rate:</b> 94.2% (industry benchmark: 87%)",
        "• <b>Average Contract Value:</b> $142,000 (12% increase from previous quarter)"
    ]
    for item in customer_list:
        story.append(Paragraph(item, styles['Normal']))
    story.append(Spacer(1, 12))
    
    # Market Position
    story.append(Paragraph("Market Position", heading_style))
    market_list = [
        "• <b>Market Share:</b> 18.5% in cloud services sector (up from 16.2%)",
        "• <b>Customer NPS Score:</b> 68 (industry average: 45)",
        "• <b>Brand Recognition:</b> 87% awareness in target market"
    ]
    for item in market_list:
        story.append(Paragraph(item, styles['Normal']))
    story.append(Spacer(1, 12))
    
    # Geographic Performance
    story.append(Paragraph("Geographic Performance", heading_style))
    geo_data = [
        ['Region', 'Revenue', 'Percentage'],
        ['North America', '$267.9M', '55.2%'],
        ['Europe', '$145.6M', '30.0%'],
        ['Asia-Pacific', '$71.7M', '14.8%'],
    ]
    geo_table = Table(geo_data, colWidths=[2*inch, 1.5*inch, 1.5*inch])
    geo_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.darkgreen),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.lightgreen),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(geo_table)
    story.append(Spacer(1, 12))
    
    # Future Outlook
    story.append(Paragraph("Future Outlook", heading_style))
    outlook_list = [
        "• <b>Q4 2024 Revenue Target:</b> $510-525 million",
        "• <b>Growth Expectation:</b> 19-22% year-over-year",
        "• <b>Operating Margin:</b> Maintain 20%+ target",
        "• <b>Strategic Initiatives:</b> Next-generation AI platform launch",
        "• <b>Expansion:</b> 5 new international markets",
        "• <b>Investment:</b> $200 million strategic acquisition budget"
    ]
    for item in outlook_list:
        story.append(Paragraph(item, styles['Normal']))
    
    # Build PDF
    doc.build(story)
    return filename

def create_academic_pptx():
    """Create academic presentation in PPTX format"""
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Machine Learning Applications in Natural Language Processing"
    subtitle.text = "A Comprehensive Study\nTransformer Architecture Evaluation\nBERT, GPT-3, and T5 Analysis"
    
    # Slide 2: Research Objectives
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Research Objectives"
    tf = content.text_frame
    tf.text = "Evaluate comparative performance of transformer models"
    
    p = tf.add_paragraph()
    p.text = "Analyze fine-tuning strategies for optimal results"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Investigate computational efficiency trade-offs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Propose optimization strategies for deployment"
    p.level = 1
    
    # Slide 3: Methodology
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Experimental Setup"
    tf = content.text_frame
    tf.text = "Hardware Configuration"
    
    p = tf.add_paragraph()
    p.text = "NVIDIA A100 GPUs (80GB VRAM)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PyTorch 1.12, Transformers 4.21, CUDA 11.7"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Training Time: 240 hours total"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Evaluation Metrics: F1 Score, BLEU, ROUGE, Accuracy"
    p.level = 1
    
    # Slide 4: Model Configurations
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Model Configurations"
    tf = content.text_frame
    tf.text = "BERT-Large: 340M parameters, 24 layers, 1024 hidden units"
    
    p = tf.add_paragraph()
    p.text = "GPT-3: 175B parameters, 96 layers, 12,288 hidden units"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "T5-Large: 770M parameters, 24 layers, 1024 hidden units"
    p.level = 0
    
    # Slide 5: Performance Results
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sentiment Analysis Results"
    tf = content.text_frame
    tf.text = "BERT-Large: 94.2% accuracy, 45ms inference"
    
    p = tf.add_paragraph()
    p.text = "T5-Large: 93.5% accuracy, 67ms inference"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "GPT-3: 91.8% accuracy, 120ms inference"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "BERT demonstrates optimal balance of performance and efficiency"
    p.level = 1
    
    # Slide 6: Key Findings
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Findings"
    tf = content.text_frame
    tf.text = "Task-Specific Performance"
    
    p = tf.add_paragraph()
    p.text = "BERT excels in understanding tasks"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "GPT-3 superior for generation tasks"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "T5 optimal for text-to-text scenarios"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Fine-tuning Impact: 8-15% performance improvement"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Hybrid approaches show 12-15% benefits"
    p.level = 0
    
    # Slide 7: Future Research
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Research Directions"
    tf = content.text_frame
    tf.text = "Investigation of newer architectures (PaLM, LaMDA, ChatGPT)"
    
    p = tf.add_paragraph()
    p.text = "Domain-specific model adaptation strategies"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Environmental impact analysis of large language models"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Development of efficient inference optimization techniques"
    p.level = 0
    
    return prs

def create_chart_image():
    """Create a performance chart as PNG"""
    # Create image
    width, height = 800, 600
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    # Try to use a better font, fall back to default if not available
    try:
        title_font = ImageFont.truetype("arial.ttf", 24)
        label_font = ImageFont.truetype("arial.ttf", 16)
        text_font = ImageFont.truetype("arial.ttf", 12)
    except:
        title_font = ImageFont.load_default()
        label_font = ImageFont.load_default()
        text_font = ImageFont.load_default()
    
    # Title
    title = "Model Performance Comparison"
    title_bbox = draw.textbbox((0, 0), title, font=title_font)
    title_width = title_bbox[2] - title_bbox[0]
    draw.text(((width - title_width) // 2, 30), title, fill='black', font=title_font)
    
    # Chart area
    chart_left = 100
    chart_right = 700
    chart_top = 100
    chart_bottom = 450
    
    # Draw axes
    draw.line([(chart_left, chart_bottom), (chart_right, chart_bottom)], fill='black', width=2)  # X-axis
    draw.line([(chart_left, chart_top), (chart_left, chart_bottom)], fill='black', width=2)  # Y-axis
    
    # Model data
    models = ['BERT-Large', 'T5-Large', 'GPT-3']
    accuracy = [94.2, 93.5, 91.8]
    colors_list = ['#1f77b4', '#ff7f0e', '#2ca02c']
    
    # Draw bars
    bar_width = 150
    bar_spacing = 170
    
    for i, (model, acc, color) in enumerate(zip(models, accuracy, colors_list)):
        x = chart_left + 50 + i * bar_spacing
        bar_height = (acc / 100) * (chart_bottom - chart_top)
        y = chart_bottom - bar_height
        
        # Draw bar
        draw.rectangle([x, y, x + bar_width, chart_bottom], fill=color, outline='black')
        
        # Draw percentage on top of bar
        percentage_text = f"{acc}%"
        text_bbox = draw.textbbox((0, 0), percentage_text, font=text_font)
        text_width = text_bbox[2] - text_bbox[0]
        draw.text((x + (bar_width - text_width) // 2, y - 25), percentage_text, fill='black', font=text_font)
        
        # Draw model name below bar
        name_bbox = draw.textbbox((0, 0), model, font=label_font)
        name_width = name_bbox[2] - name_bbox[0]
        draw.text((x + (bar_width - name_width) // 2, chart_bottom + 10), model, fill='black', font=label_font)
    
    # Y-axis labels
    for i in range(0, 101, 10):
        y = chart_bottom - (i / 100) * (chart_bottom - chart_top)
        draw.text((chart_left - 40, y - 10), f"{i}%", fill='black', font=text_font)
        # Grid lines
        draw.line([(chart_left, y), (chart_right, y)], fill='lightgray', width=1)
    
    # X-axis label
    xlabel = "Models"
    xlabel_bbox = draw.textbbox((0, 0), xlabel, font=label_font)
    xlabel_width = xlabel_bbox[2] - xlabel_bbox[0]
    draw.text(((width - xlabel_width) // 2, chart_bottom + 50), xlabel, fill='black', font=label_font)
    
    # Y-axis label (rotated effect with multiple lines)
    draw.text((20, height // 2 - 40), "A", fill='black', font=label_font)
    draw.text((20, height // 2 - 20), "c", fill='black', font=label_font)
    draw.text((20, height // 2), "c", fill='black', font=label_font)
    draw.text((20, height // 2 + 20), "u", fill='black', font=label_font)
    draw.text((20, height // 2 + 40), "r", fill='black', font=label_font)
    draw.text((20, height // 2 + 60), "a", fill='black', font=label_font)
    draw.text((20, height // 2 + 80), "c", fill='black', font=label_font)
    draw.text((20, height // 2 + 100), "y", fill='black', font=label_font)
    
    return img

def main():
    """Generate all document types"""
    base_path = Path(__file__).parent / "../documents"
    
    print("Creating diverse document formats...")
    
    # Create DOCX (Technical)
    print("Creating technical documentation (DOCX)...")
    docx_doc = create_technical_docx()
    docx_path = base_path / "technical" / "api_documentation.docx"
    docx_doc.save(str(docx_path))
    print(f"Created: {docx_path}")
    
    # Create PDF (Business)
    print("Creating business report (PDF)...")
    pdf_path = create_business_pdf()
    print(f"Created: {pdf_path}")
    
    # Create PPTX (Academic)
    print("Creating academic presentation (PPTX)...")
    pptx_doc = create_academic_pptx()
    pptx_path = base_path / "academic" / "ml_nlp_research_presentation.pptx"
    pptx_doc.save(str(pptx_path))
    print(f"Created: {pptx_path}")
    
    # Create PNG chart
    print("Creating performance chart (PNG)...")
    chart_img = create_chart_image()
    chart_path = base_path / "academic" / "performance_comparison.png"
    chart_img.save(str(chart_path))
    print(f"Created: {chart_path}")
    
    # Keep the existing CSV file
    print("Structured data (CSV) already exists")
    
    print("\nAll document formats created successfully!")
    print("Document types now include: DOCX, PDF, PPTX, PNG, CSV")

if __name__ == "__main__":
    main()